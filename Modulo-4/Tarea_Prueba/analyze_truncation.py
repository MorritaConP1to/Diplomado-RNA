import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Pt

prs = Presentation(r'C:\Users\Darck\Downloads\Tarea\Backpropagation_from_Scratch.pptx')

slide_w = 10.0
slide_h = 5.625

for idx, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        w = shape.width / 914400
        h = shape.height / 914400
        left = shape.left / 914400
        top = shape.top / 914400
        
        tf = shape.text_frame
        full_text = tf.text
        
        if not full_text.strip():
            continue
        
        # Check if text ends with "..." or "..." in the markitdown output
        # The markitdown output showed truncated text (ending with "...")
        
        # Estimate line count needed
        total_chars = len(full_text)
        avg_font_size = 12
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size:
                    avg_font_size = run.font.size.pt
                    break
        
        # Approximate char per inch based on font size (roughly pt/2 chars per inch)
        chars_per_inch = 72 / avg_font_size * 1.8
        chars_per_line = w * chars_per_inch
        total_lines = total_chars / max(chars_per_line, 1) if chars_per_line > 0 else 999
        
        # Approximate line height (1.2 * font size)
        line_height_in = avg_font_size * 1.2 / 72
        needed_height = total_lines * line_height_in
        
        # Check truncation at bottom
        right_edge = left + w
        bottom_edge = top + h
        
        trunc_issues = []
        if bottom_edge > slide_h:
            trunc_issues.append(f'EXTENDS_BEYOND_SLIDE_BOTTOM by {bottom_edge - slide_h:.2f}"')
        if right_edge > slide_w:
            trunc_issues.append(f'EXTENDS_BEYOND_SLIDE_RIGHT by {right_edge - slide_w:.2f}"')
        
        # Check if text ends with partial content (last visible line)
        last_line = full_text.split('\n')[-1].strip()
        if last_line and len(last_line) > 5:
            pass  # This is normal
        
        # Check if text box is too small for content
        if needed_height > h * 1.1:
            trunc_issues.append(f'POSSIBLE_CLIP: ~{total_lines:.0f} lines need ~{needed_height:.1f}" but box is {h:.1f}"')
        
        if trunc_issues:
            print(f'Slide {idx} [{shape.name}]: {" | ".join(trunc_issues)}')
            print(f'  text="{full_text[:100]}"')

# Check for "..." in text content (truncation markers)
print()
print('=== TEXT THAT MAY BE TRUNCATED (checking XXX pattern) ===')
for idx, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text
        # Check for placeholder text
        if any(p in t.lower() for p in ['xxxx', 'lorem', 'ipsum', 'placeholder']):
            print(f'Slide {idx} [{shape.name}]: PLACEHOLDER TEXT FOUND: {t[:100]}')
