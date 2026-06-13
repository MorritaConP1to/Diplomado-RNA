import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn

prs = Presentation(r'C:\Users\Darck\Downloads\Tarea\Backpropagation_from_Scratch.pptx')

slide_w = 10.0
slide_h = 5.625

# Check if shapes have fill colors (for the decorative circles)
print('=== DECORATIVE SHAPES ===')
for idx, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            if shape.width/914400 > 1.5:  # Large non-text shapes
                left = shape.left/914400
                top = shape.top/914400
                w = shape.width/914400
                h = shape.height/914400
                print(f'Slide {idx}: Large non-text shape "{shape.name}" @({left:.1f},{top:.1f}) size({w:.1f}x{h:.1f})')

print()
print('=== OVERLAP ANALYSIS (specific pairs) ===')

# Slide 1: title and subtitle overlap
s1_title = None
s1_sub = None
for shape in prs.slides[0].shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text
        if 'Backpropagation' in t:
            s1_title = shape
        if 'The algorithm' in t:
            s1_sub = shape

if s1_title and s1_sub:
    t_bot = (s1_title.top + s1_title.height) / 914400
    s_top = s1_sub.top / 914400
    print(f'Slide 1: Title bottom={t_bot:.2f}, Subtitle top={s_top:.2f}, gap={s_top - t_bot:.2f}"')
    if t_bot > s_top:
        print(f'  *** OVERLAP: Title and subtitle overlap by {t_bot - s_top:.2f}"')

# Slide 14: table rows overlap code box
slide14 = prs.slides[13]
for shape in slide14.shapes:
    if shape.has_text_frame and ('[1, 1]' in shape.text_frame.text):
        t_bot = (shape.top + shape.height) / 914400
        print(f'Slide 14: "[1,1]" row bottom={t_bot:.2f}"')

# Shape 16 (code box) on slide 14
for shape in slide14.shapes:
    if shape.name == 'Shape 16':
        cb_top = shape.top / 914400
        print(f'Slide 14: Code box top={cb_top:.2f}"')
        if cb_top < t_bot:
            print(f'  *** OVERLAP: Data row extends into code box by {t_bot - cb_top:.2f}"')

print()
print('=== TEXT BOX INTERNAL PADDING CHECK ===')
for idx, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        # Check internal margins
        try:
            l = tf.margin_left
            r = tf.margin_right
            t = tf.margin_top
            b = tf.margin_bottom
            if l is not None or r is not None or t is not None or b is not None:
                vals = []
                if l is not None: vals.append(f'L={l/914400:.2f}"')
                if r is not None: vals.append(f'R={r/914400:.2f}"')
                if t is not None: vals.append(f'T={t/914400:.2f}"')
                if b is not None: vals.append(f'B={b/914400:.2f}"')
                if vals:
                    print(f'Slide {idx} [{shape.name}]: margins=[{", ".join(vals)}]')
        except:
            pass
