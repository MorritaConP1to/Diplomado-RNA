import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

prs = Presentation(r'C:\Users\Darck\Downloads\Tarea\Backpropagation_from_Scratch.pptx')

def lum(r, g, b):
    return 0.299*r + 0.587*g + 0.114*b

for idx, slide in enumerate(prs.slides, 1):
    print(f'=== SLIDE {idx} COLOR DETAIL ===')
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                info = []
                if run.font.size:
                    info.append(f'{run.font.size.pt}pt')
                if run.font.color and run.font.color.type is not None:
                    try:
                        c = run.font.color.rgb
                        r, g, b = c[0], c[1], c[2]
                        l = lum(r, g, b)
                        info.append(f'RGB({r},{g},{b}) lum={l:.0f}')
                        # Low contrast if too light/too dark
                        if 120 < l < 200:
                            info.append('LOW_CONTRAST!')
                    except:
                        info.append('theme_color')
                t = run.text[:60].replace('\n','|')
                if info:
                    print(f'  [{shape.name}] {", ".join(info)} text="{t}"')
    print()

# Check slide fill colors
print('=== SLIDE BACKGROUNDS ===')
for idx, slide in enumerate(prs.slides, 1):
    try:
        bg = slide.background.fill
        if bg.type == 1:
            c = bg.fore_color.rgb
            r, g, b = c[0], c[1], c[2]
            print(f'Slide {idx}: BG=RGB({r},{g},{b}) lum={lum(r,g,b):.0f}')
        else:
            print(f'Slide {idx}: BG type={bg.type}')
    except:
        print(f'Slide {idx}: BG unknown')
