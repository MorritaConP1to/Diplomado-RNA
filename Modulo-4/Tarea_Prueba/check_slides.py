import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

prs = Presentation(r'C:\Users\Darck\Downloads\Tarea\Backpropagation_from_Scratch.pptx')
sw = prs.slide_width
sh = prs.slide_height
print(f'Slide dimensions: {sw/914400:.2f}" x {sh/914400:.2f}"')
print('=' * 80)

targets = [1, 3, 4, 13, 14, 15, 18]
for n, slide in enumerate(prs.slides, 1):
    if n not in targets:
        continue
    print(f'\n--- Slide {n} ({len(list(slide.shapes))} shapes) ---')
    for i, shape in enumerate(slide.shapes):
        left = shape.left
        top = shape.top
        w = shape.width
        h = shape.height
        right = left + w
        bottom = top + h
        text = ''
        if shape.has_text_frame:
            text = shape.text_frame.text[:100].replace('\n', ' | ')
        print(f'  [{i:2d}] L={left/914400:.3f}" T={top/914400:.3f}" W={w/914400:.3f}" H={h/914400:.3f}" R={right/914400:.3f}" B={bottom/914400:.3f}"  txt="{text}"')
