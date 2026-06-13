import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation(r'C:\Users\Darck\Downloads\Tarea\Backpropagation_from_Scratch.pptx')

def emu_to_inches(emu):
    if emu is None: return None
    return round(emu / 914400, 3)

slide_w = emu_to_inches(prs.slide_width)
slide_h = emu_to_inches(prs.slide_height)
print(f'Slide dimensions: {slide_w}" x {slide_h}"')
print()

for idx, slide in enumerate(prs.slides, 1):
    print(f'=== SLIDE {idx} ===')
    layout_name = slide.slide_layout.name if slide.slide_layout else 'None'
    print(f'Layout: {layout_name}')
    
    shapes_info = []
    for shape in slide.shapes:
        sname = shape.name or ''
        left = emu_to_inches(shape.left)
        top = emu_to_inches(shape.top)
        w = emu_to_inches(shape.width)
        h = emu_to_inches(shape.height)
        right = round(left + w, 3) if left is not None and w is not None else None
        bottom = round(top + h, 3) if top is not None and h is not None else None
        
        issues = []
        
        # Check margins from slide edges
        if left is not None and left < 0.5 and w is not None and w < slide_w * 0.8:
            issues.append(f'LEFT_MARGIN={left}" < 0.5"')
        if top is not None and top < 0.3:
            issues.append(f'TOP_MARGIN={top}" < 0.3"')
        if right is not None and right > slide_w - 0.4:
            issues.append(f'RIGHT_EDGE={right}" (max {slide_w - 0.4}") < 0.4" from edge')
        if bottom is not None and bottom > slide_h - 0.3:
            issues.append(f'BOTTOM_EDGE={bottom}" (max {slide_h - 0.3}") < 0.3" from edge')
        
        # Check for very narrow text boxes
        if w is not None and w < 1.2 and shape.has_text_frame and shape.text_frame.text.strip():
            issues.append(f'NARROW_BOX={w}" wide')
        
        shape_type = str(shape.shape_type)
        text_preview = ''
        if shape.has_text_frame:
            txt = shape.text_frame.text[:100].replace('\n', '|')
            text_preview = f' text="{txt}"'
            # Check font sizes and colors
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        pt = run.font.size.pt
                        if w and pt > 0:
                            est_chars_per_line = int(w * 72 / pt * 0.55)
                            if est_chars_per_line < 12 and len(run.text) > 20 and est_chars_per_line > 0:
                                issues.append(f'OVERFLOW_RISK: {pt}pt in {w}" box ~{est_chars_per_line}ch/line')
                    if run.font.color and run.font.color.rgb:
                        try:
                            r, g, b = run.font.color.rgb[0], run.font.color.rgb[1], run.font.color.rgb[2]
                            # Check for low contrast (light gray on light bg)
                            lum = 0.299*r + 0.587*g + 0.114*b
                            if 180 < lum < 240:
                                issues.append(f'LOW_CONTRAST: RGB({r},{g},{b}) luminance={lum:.0f}')
                        except:
                            pass
        
        shapes_info.append({
            'name': sname,
            'type': shape_type,
            'left': left, 'top': top, 'w': w, 'h': h,
            'right': right, 'bottom': bottom,
            'text': shape.text_frame.text if shape.has_text_frame else '',
            'issues': issues
        })
        
        issue_str = ' *** ' + '; '.join(issues) if issues else ''
        line = f'  [{shape_type}] {sname} @({left},{top}) size({w}x{h}){text_preview}{issue_str}'
        print(line.encode('utf-8', 'replace').decode('utf-8', 'replace'))
    
    # Check for overlapping elements
    for i in range(len(shapes_info)):
        for j in range(i+1, len(shapes_info)):
            a = shapes_info[i]
            b = shapes_info[j]
            if None in (a['left'], a['right'], b['left'], b['right']):
                continue
            # Check overlap
            x_overlap = max(0, min(a['right'], b['right']) - max(a['left'], b['left']))
            y_overlap = max(0, min(a['bottom'], b['bottom']) - max(a['top'], b['top']))
            if x_overlap > 0 and y_overlap > 0:
                overlap_area = x_overlap * y_overlap
                # Only report if significant overlap
                if overlap_area > 0.5 and a['text'].strip() and b['text'].strip():
                    print(f'    OVERLAP: "{a["name"]}" and "{b["name"]}" overlap by {x_overlap:.2f}x{y_overlap:.2f}"')
    
    # Check gaps between consecutive elements
    for i in range(len(shapes_info)):
        for j in range(len(shapes_info)):
            if i == j: continue
            a = shapes_info[i]
            b = shapes_info[j]
            if None in (a['left'], a['right'], a['top'], a['bottom'], b['left'], b['right'], b['top'], b['bottom']):
                continue
            # Horizontal gap: a is left of b
            if a['right'] <= b['left']:
                h_gap = b['left'] - a['right']
                if 0 < h_gap < 0.3:
                    print(f'    TIGHT_H_GAP: {h_gap:.2f}" between "{a["name"]}" and "{b["name"]}"')
            # Vertical gap: a is above b
            if a['bottom'] <= b['top']:
                v_gap = b['top'] - a['bottom']
                if 0 < v_gap < 0.15:
                    print(f'    TIGHT_V_GAP: {v_gap:.2f}" between "{a["name"]}" and "{b["name"]}"')
    
    print()

# Check background colors
print('=== BACKGROUND COLORS ===')
for idx, slide in enumerate(prs.slides, 1):
    bg = slide.background
    if bg.fill.type is not None:
        try:
            fill = bg.fill
            if fill.type == 1:  # Solid
                print(f'Slide {idx}: solid fill')
        except:
            print(f'Slide {idx}: fill type={bg.fill.type}')
    else:
        print(f'Slide {idx}: no explicit fill (inherited)')
